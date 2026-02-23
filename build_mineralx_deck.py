#!/usr/bin/env python3
"""
MineralX Strategic Fund Vision Deck - World Class Version
Clean institutional design: White/Charcoal/Slate (NOT gold+navy)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 widescreen layout
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette - Institutional Modern (NOT gold+navy)
COLORS = {
    'white': RGBColor(255, 255, 255),
    'charcoal': RGBColor(45, 45, 50),
    'dark_slate': RGBColor(30, 35, 45),
    'slate': RGBColor(100, 110, 125),
    'light_gray': RGBColor(240, 242, 245),
    'medium_gray': RGBColor(180, 185, 195),
    'accent_teal': RGBColor(0, 130, 140),
    'accent_amber': RGBColor(210, 140, 40),
}

def add_title_slide(prs, title, subtitle=""):
    """Title slide with clean minimal design"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    # Dark accent bar at top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['charcoal']
    top_bar.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    p.font.name = "Helvetica Neue"
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['slate']
        p.font.name = "Helvetica Neue"
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(6), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL | FEBRUARY 2025"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['medium_gray']
    p.font.name = "Helvetica Neue"
    
    return slide

def add_content_slide(prs, title, subtitle=""):
    """Standard content slide with header"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    # Top accent line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLORS['accent_teal']
    top_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    p.font.name = "Helvetica Neue"
    
    # Subtitle if provided
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['slate']
        p.font.name = "Helvetica Neue"
    
    return slide, Inches(2)

# ============================================
# BUILD THE DECK
# ============================================

print("Building MineralX Strategic Fund Vision Deck...")

# Slide 1: Title
add_title_slide(prs, "MINERALX", "Integrated Mining & Processing Platform\nStrategic Fund Vision Deck")

# Slide 2: Why Now
slide, y = add_content_slide(prs, "Why Now", "Precious metals demand is surging. Critical minerals are strategic infrastructure.")

# Four data boxes
boxes = [
    ("Structural Demand", ["Gold: $2,900+/oz — 30-year high", "Critical minerals: $770B by 2030"], COLORS['accent_teal']),
    ("Policy Acceleration", ["Australia: $4B committed", "US IRA supply chain incentives"], COLORS['accent_amber']),
    ("Processing Bottleneck", ["<5% of rare earth processing", "70%+ gold capacity constrained"], COLORS['accent_teal']),
    ("Fragmented Opportunity", ["1,000+ dormant QLD tenements", "No integrated platforms exist"], COLORS['slate'])
]

box_width = Inches(5.8)
box_height = Inches(2.2)
for i, (title, points, color) in enumerate(boxes):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.8)
    y_pos = Inches(2.2) if i < 2 else Inches(4.6)
    
    # Colored accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, Inches(0.08), box_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()
    
    # Title
    t_box = slide.shapes.add_textbox(x + Inches(0.2), y_pos + Inches(0.2), box_width - Inches(0.4), Inches(0.4))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    p.font.name = "Helvetica Neue"
    
    # Points
    p_box = slide.shapes.add_textbox(x + Inches(0.2), y_pos + Inches(0.7), box_width - Inches(0.4), box_height - Inches(0.9))
    tf = p_box.text_frame
    tf.word_wrap = True
    for j, pt in enumerate(points):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['slate']
        p.font.name = "Helvetica Neue"
        p.space_before = Pt(4)

# Slide 3: Industry Problem
slide, y = add_content_slide(prs, "The Industry Problem", "Three structural failures creating an opportunity")

cols = [
    ("Processing Gap", ["Infrastructure bottlenecks", "$50M+ typical plant capex", "80% of juniors never produce"], COLORS['accent_teal']),
    ("Capital Constraint", ["Explorers lack funding", "Ore sits waiting for offtake", "Dilution destroys value"], COLORS['accent_amber']),
    ("Fragmentation", ["No integrated platforms", "Toll treatment = margin loss", "Stranded assets"], COLORS['slate'])
]

col_width = Inches(3.8)
for i, (title, points, color) in enumerate(cols):
    x = Inches(0.8) + (col_width + Inches(0.3)) * i
    
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    # Title
    t_box = slide.shapes.add_textbox(x, y + Inches(0.2), col_width, Inches(0.4))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    
    # Points
    p_box = slide.shapes.add_textbox(x, y + Inches(0.7), col_width, Inches(3))
    tf = p_box.text_frame
    tf.word_wrap = True
    for j, pt in enumerate(points):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['slate']
        p.space_before = Pt(6)

# Slide 4: The MineralX Platform (with table)
slide, y = add_content_slide(prs, "The MineralX Platform", "End-to-end mining and processing")

# Simple comparison
headers = ["", "Traditional Junior", "MineralX"]
rows = [
    ["Timeline to Revenue", "7-10 years", "12-18 months"],
    ["Capex Required", "$50M–$150M+", "$5M–$15M modular"],
    ["Processing Control", "Third-party toll", "Owned infrastructure"],
    ["Margin Profile", "30-40%", "60-70% integrated"],
]

shape = slide.shapes.add_table(5, 3, Inches(1), Inches(2.3), Inches(11.3), Inches(3.5))
table = shape.table

for i, row_data in enumerate([headers] + rows):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['charcoal']
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = COLORS['white']
                paragraph.font.bold = True
        elif j == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['light_gray']
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
        elif j == 2:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = COLORS['accent_teal']
                paragraph.font.bold = True

# Slide 5: Traction
slide, y = add_content_slide(prs, "Traction & Proof Points", "ML 30139: Active gold production")

metrics = [
    ("Current Asset", "ML 30139", "Active Mining Lease, Queensland"),
    ("Processing", "Pilot Operations", "Gravity + CIL optimisation"),
    ("Timeline", "12-18 Months", "To commercial revenue"),
    ("Pipeline", "Multiple EPMs", "Under LOI / review")
]

for i, (label, value, detail) in enumerate(metrics):
    y_pos = Inches(2.3) + i * Inches(1.1)
    
    v_box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(4), Inches(0.5))
    tf = v_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_teal']
    
    l_box = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.45), Inches(3), Inches(0.3))
    tf = l_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['slate']
    
    d_box = slide.shapes.add_textbox(Inches(5), y_pos, Inches(6), Inches(0.4))
    tf = d_box.text_frame
    p = tf.paragraphs[0]
    p.text = detail
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['charcoal']

# Slide 6: Three Pillars
slide, y = add_content_slide(prs, "Three Integrated Engines", "Platform architecture")

pillars = [
    ("Mining & Exploration", ["ML 30139 gold production", "EPM consolidation", "12-18 months to cash flow"], COLORS['accent_teal']),
    ("Processing R&D", ["Gravity + CIL optimisation", "Modular plant design", "Recovery improvements"], COLORS['accent_amber']),
    ("Manufacturing", ["Modular processing units", "Lower capex vs industry", "Scalable deployment"], COLORS['slate'])
]

for i, (title, points, color) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(4.2)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    t_box = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(3.8), Inches(0.4))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    
    p_box = slide.shapes.add_textbox(x, y + Inches(0.7), Inches(3.8), Inches(2.5))
    tf = p_box.text_frame
    tf.word_wrap = True
    for j, pt in enumerate(points):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['slate']
        p.space_before = Pt(6)

# Slide 7: Capital Structure
slide, y = add_content_slide(prs, "Capital Structure", "Clean ownership with strategic flexibility")

cap_data = [
    ("Founder Holdings (Economic)", "65M shares", "65%", "50%", COLORS['accent_teal']),
    ("Founder SPV (Control)", "5M shares", "5%", "3.85%", COLORS['accent_teal']),
    ("Strategic Anchor Partner", "5M shares", "5%", "3.85%", COLORS['accent_amber']),
    ("ESOP Pool", "10M shares", "10%", "7.69%", COLORS['slate']),
    ("Future Capital Reserve", "20M shares", "15%", "15.38%", COLORS['light_gray']),
]

headers = ["Shareholder", "Shares", "Pre-Dilution", "Fully Diluted"]
x_positions = [Inches(0.8), Inches(7), Inches(9), Inches(11)]

for i, h in enumerate(headers):
    box = slide.shapes.add_textbox(x_positions[i], Inches(2.3), Inches(2.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['slate']

for i, (name, shares, pre, diluted, color) in enumerate(cap_data):
    y_pos = Inches(2.8) + i * Inches(0.6)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Inches(0.12), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    slide.shapes.add_textbox(Inches(1.1), y_pos, Inches(5.5), Inches(0.5)).text_frame.paragraphs[0].text = name
    slide.shapes.add_textbox(Inches(7), y_pos, Inches(1.8), Inches(0.5)).text_frame.paragraphs[0].text = shares
    slide.shapes.add_textbox(Inches(9), y_pos, Inches(1.5), Inches(0.5)).text_frame.paragraphs[0].text = pre
    slide.shapes.add_textbox(Inches(11), y_pos, Inches(1.5), Inches(0.5)).text_frame.paragraphs[0].text = diluted

# Slide 8: THE ASK
slide, y = add_content_slide(prs, "The Investment Opportunity", "Strategic anchor partner for platform-scale mining")

ask_metrics = [
    ("Raise", "A$10,000,000", "Two-tranche structure"),
    ("Valuation", "A$200,000,000", "Post-money"),
    ("Equity", "5.0%", "5,000,000 shares"),
    ("Price", "A$2.00", "Per share")
]

for i, (label, value, detail) in enumerate(ask_metrics):
    x = Inches(0.8) + (i % 2) * Inches(6)
    y_pos = Inches(2.3) + (i // 2) * Inches(2.2)
    
    v_box = slide.shapes.add_textbox(x, y_pos, Inches(5.5), Inches(0.8))
    tf = v_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_teal']
    
    l_box = slide.shapes.add_textbox(x, y_pos + Inches(0.8), Inches(5.5), Inches(0.4))
    tf = l_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    
    d_box = slide.shapes.add_textbox(x, y_pos + Inches(1.2), Inches(5.5), Inches(0.4))
    tf = d_box.text_frame
    p = tf.paragraphs[0]
    p.text = detail
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['slate']

# Slide 9: Use of Proceeds
slide, y = add_content_slide(prs, "Use of Proceeds", "Tranche A: $6M on signing | Tranche B: $4M on milestone")

proceeds = [
    ("Processing Plant Capex", "35%", "$3.5M", COLORS['accent_teal']),
    ("Mining Operations", "30%", "$3.0M", COLORS['accent_amber']),
    ("Acquisition Pipeline", "20%", "$2.0M", COLORS['accent_teal']),
    ("Working Capital", "10%", "$1.0M", COLORS['slate']),
    ("Corporate", "5%", "$0.5M", COLORS['light_gray']),
]

for i, (category, pct, amount, color) in enumerate(proceeds):
    y_pos = Inches(2.3) + i * Inches(0.85)
    pct_val = float(pct.replace('%', '')) / 100
    
    # Bar background
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos + Inches(0.3), Inches(4), Inches(0.22))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = COLORS['light_gray']
    bar_bg.line.fill.background()
    
    # Bar fill
    bar_fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos + Inches(0.3), Inches(4 * pct_val), Inches(0.22))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = color
    bar_fill.line.fill.background()
    
    # Labels
    slide.shapes.add_textbox(Inches(5), y_pos, Inches(3.5), Inches(0.35)).text_frame.paragraphs[0].text = category
    slide.shapes.add_textbox(Inches(5), y_pos + Inches(0.35), Inches(3), Inches(0.3)).text_frame.paragraphs[0].text = f"{pct} ({amount})"

# Slide 10: Team
slide, y = add_content_slide(prs, "Leadership", "Founder-operator with technical and commercial track record")

team = [
    ("Ayrton Mansi", "Founder & CEO", [
        "Chemical engineering background",
        "Hands-on mining & processing experience", 
        "Property development track record",
        "ML 30139 operator"
    ])
]

for name, role, points in team:
    # Name
    n_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(6), Inches(0.6))
    tf = n_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['charcoal']
    
    # Role
    r_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.6), Inches(6), Inches(0.4))
    tf = r_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['accent_teal']
    
    # Points
    p_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(1.1), Inches(11), Inches(3))
    tf = p_box.text_frame
    tf.word_wrap = True
    for j, pt in enumerate(points):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['slate']
        p.space_before = Pt(8)

# Slide 11: Deployment Timeline
slide, y = add_content_slide(prs, "Deployment Timeline", "Phased execution from pilot to platform scale")

phases = [
    ("PHASE 1: PILOT", "Q1–Q2 2026", ["Pilot processing operational", "First gold pour", "Toll treatment agreements"], COLORS['accent_teal']),
    ("PHASE 2: SCALE", "Q3 2026–Q2 2027", ["Commercial production ramp", "EPM acquisitions", "Processing expansion"], COLORS['accent_amber']),
    ("PHASE 3: PLATFORM", "2027+", ["Regional hub model", "Multiple commodities", "Institutional capital"], COLORS['slate'])
]

for i, (phase, timing, items, color) in enumerate(phases):
    x = Inches(0.8) + i * Inches(4.1)
    
    # Phase header
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    p_box = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(3.8), Inches(0.5))
    tf = p_box.text_frame
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    
    t_box = slide.shapes.add_textbox(x, y + Inches(0.5), Inches(3.8), Inches(0.4))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = timing
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['slate']
    
    items_box = slide.shapes.add_textbox(x, y + Inches(1), Inches(3.8), Inches(3))
    tf = items_box.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['charcoal']
        p.space_before = Pt(6)

# Slide 12: Closing
slide, y = add_content_slide(prs, "The Opportunity", "First-mover platform in fragmented Australian mining")

closing_points = [
    "A$10M raise at A$200M post-money — 5% strategic equity",
    "Active production asset: ML 30139 with clear expansion path",
    "12-18 months to commercial revenue vs. 7-10 years industry standard",
    "Integrated model: 60-70% margins vs. 30-40% toll treatment",
    "Consolidation opportunity: 1,000+ dormant tenements in target region"
]

for i, point in enumerate(closing_points):
    y_pos = y + i * Inches(0.8)
    
    # Bullet
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y_pos + Inches(0.15), Inches(0.12), Inches(0.12))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = COLORS['accent_teal']
    bullet.line.fill.background()
    
    # Text
    p_box = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(11), Inches(0.6))
    tf = p_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(15)
    p.font.color.rgb = COLORS['charcoal']

# Contact
contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(12), Inches(0.8))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "Ayrton Mansi | ayrton@mineral-x.com.au"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['accent_teal']

# Save the presentation
output_path = "/Users/ayrtonmansi/.openclaw/workspace/MineralX_Strategic_Fund_Vision_Deck.pptx"
prs.save(output_path)
print(f"✓ Deck created: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
