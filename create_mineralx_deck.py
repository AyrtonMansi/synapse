#!/usr/bin/env python3
"""
MineralX Strategic Fund Vision Deck Builder
Creates a professional 18-slide PPTX presentation for strategic capital partners
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Color scheme
DARK_NAVY = RGBColor(31, 41, 55)    # #1F2937
GOLD = RGBColor(212, 168, 67)       # #D4A843
WHITE = RGBColor(255, 255, 255)     # #FFFFFF
LIGHT_GRAY = RGBColor(248, 250, 252) # #F8FAFC

def create_blank_presentation():
    """Create a blank presentation with 16:9 aspect ratio"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # MINERALX title (large)
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    title_frame = title.text_frame
    title_frame.text = "MINERALX"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(0.8))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Integrated Mining & Processing Platform"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.8))
    footer_frame = footer.text_frame
    footer_frame.text = "Strategic Fund Vision Deck | Confidential | February 2026"
    p = footer_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER

def add_why_now_slide(prs):
    """Slide 2: Why Now"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "Why Now"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.333), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Structural demand meets infrastructure gap"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.LEFT
    
    # Four boxes
    boxes = [
        ("Structural Demand", "Gold at $2,900+/oz (30-year highs). Critical minerals market projected $770B by 2030. Multi-decade demand from electrification, AI infrastructure, energy storage."),
        ("Policy Acceleration", "Australia's Critical Minerals Strategy: $4B+ committed. US IRA incentives. OECD supply chain security mandates prioritise domestic processing."),
        ("Processing Bottleneck", "Mining output constrained by processing infrastructure, not geology. <5% of global REE processing capacity in Australia."),
        ("Fragmented Opportunity", "1,000+ dormant gold tenements in QLD. Fragmented ownership, no processing access. Consolidation-ready.")
    ]
    
    x_positions = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]
    y_pos = Inches(2.2)
    box_width = Inches(2.8)
    box_height = Inches(4.5)
    
    for i, (title_text, content_text) in enumerate(boxes):
        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y_pos, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = GOLD
        box.line.width = Pt(2)
        
        # Box title
        box_title = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), y_pos + Inches(0.2), box_width - Inches(0.2), Inches(0.6))
        box_title_frame = box_title.text_frame
        box_title_frame.text = title_text
        p = box_title_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        # Box content
        box_content = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), y_pos + Inches(0.9), box_width - Inches(0.2), box_height - Inches(1.1))
        box_content_frame = box_content.text_frame
        box_content_frame.text = content_text
        box_content_frame.word_wrap = True
        p = box_content_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.LEFT

def add_industry_problem_slide(prs):
    """Slide 3: Industry Problem (Dark background)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "Value is Structurally Trapped"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Banner text
    banner = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2))
    banner_frame = banner.text_frame
    banner_frame.text = "The global race for critical minerals is driven by structural demand from electrification, energy storage, and strategic supply chain policy — creating multi-decade growth tailwinds for companies that can deliver integrated processing, consolidation, and scalable production."
    banner_frame.word_wrap = True
    p = banner_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Three problem cards
    problems = [
        ("Fragmented Ownership", "Mineral districts divided across multiple small operators with limited capital. Economies of scale remain unrealised."),
        ("Long Lead Times", "Traditional development requires 5-10 years to first production. Capital inefficiently allocated based on speculation."),
        ("No Processing Access", "Junior miners lack capital for processing infrastructure. Ore sits in the ground waiting for offtake.")
    ]
    
    x_positions = [Inches(1), Inches(4.7), Inches(8.4)]
    y_pos = Inches(3.2)
    box_width = Inches(3.5)
    box_height = Inches(3.5)
    
    for i, (title_text, content_text) in enumerate(problems):
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], y_pos, box_width, box_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GOLD
        card.line.width = Pt(3)
        
        # Card title
        card_title = slide.shapes.add_textbox(x_positions[i] + Inches(0.2), y_pos + Inches(0.2), box_width - Inches(0.4), Inches(0.6))
        card_title_frame = card_title.text_frame
        card_title_frame.text = title_text
        p = card_title_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        # Card content
        card_content = slide.shapes.add_textbox(x_positions[i] + Inches(0.2), y_pos + Inches(1), box_width - Inches(0.4), box_height - Inches(1.2))
        card_content_frame = card_content.text_frame
        card_content_frame.text = content_text
        card_content_frame.word_wrap = True
        p = card_content_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.LEFT

def add_platform_slide(prs):
    """Slide 4: The MineralX Platform"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(0.7))
    title_frame = title.text_frame
    title_frame.text = "The MineralX Platform"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Processing infrastructure to consolidate fragmented mineral districts"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Comparison table
    table_data = [
        ["", "MineralX", "Traditional Junior"],
        ["Time to Revenue", "12-18 months", "5-10 years"],
        ["Processing Access", "Pilot + toll treatment", "Offtake dependent"],
        ["Capital Risk", "Modular, incremental", "Binary, monolithic"],
        ["Scale Strategy", "Regional consolidation", "Single asset focus"],
        ["Cash Flow Model", "Parallel mining + exploration", "Sequential development"]
    ]
    
    # Create table manually with shapes
    start_x = Inches(2)
    start_y = Inches(2)
    col_width = Inches(3)
    row_height = Inches(0.6)
    
    for row_idx, row in enumerate(table_data):
        for col_idx, cell_text in enumerate(row):
            x = start_x + (col_idx * col_width)
            y = start_y + (row_idx * row_height)
            
            # Cell background
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width, row_height)
            if row_idx == 0:  # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            elif col_idx == 0:  # Row labels
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            elif col_idx == 1:  # MineralX column
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD
            else:  # Traditional column
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            
            cell.line.color.rgb = DARK_NAVY
            cell.line.width = Pt(1)
            
            # Cell text
            if cell_text:  # Skip empty cells
                cell_textbox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), col_width - Inches(0.2), row_height - Inches(0.1))
                cell_textbox_frame = cell_textbox.text_frame
                cell_textbox_frame.text = cell_text
                p = cell_textbox_frame.paragraphs[0]
                p.font.name = "Calibri"
                if row_idx == 0:  # Header
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.size = Pt(12)
                    if col_idx == 0:
                        p.font.bold = True
                        p.font.color.rgb = DARK_NAVY
                    elif col_idx == 1:
                        p.font.color.rgb = WHITE
                        p.font.bold = True
                    else:
                        p.font.color.rgb = DARK_NAVY
                p.alignment = PP_ALIGN.CENTER
    
    # Summary
    summary = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(1))
    summary_frame = summary.text_frame
    summary_frame.text = "MineralX consolidates fragmented Australian mineral districts through processing access, generating near-term cash flow while building long-term district-scale positions."
    summary_frame.word_wrap = True
    p = summary_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER

def add_traction_slide(prs):
    """Slide 5: Traction & Proof Points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "Current Portfolio"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Traction points
    traction_points = [
        "• ML 30139: Active gold mining lease, Queensland",
        "• Pilot processing: Gravity concentration operational", 
        "• Toll treatment: Third-party processing agreements in pipeline",
        "• Assay data confirming grades",
        "• Processing expertise: Chemical engineering + hands-on metallurgical experience",
        "• Multiple EPMs under consolidation review"
    ]
    
    points_text = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.333), Inches(4))
    points_frame = points_text.text_frame
    points_frame.text = "\n\n".join(traction_points)
    points_frame.word_wrap = True
    
    for i, p in enumerate(points_frame.paragraphs):
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.LEFT

def add_growth_strategy_slide(prs):
    """Slide 6: Growth Strategy"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(0.7))
    title_frame = title.text_frame
    title_frame.text = "Growth Strategy"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Systematic consolidation of underexplored mining districts"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Three columns
    columns = [
        ("Underexplored Areas", "• Historic mining districts\n• Limited modern exploration\n• Infrastructure constraints\n• Fragmented tenement ownership"),
        ("Historic Mining Assets", "• Known mineralisation\n• Existing infrastructure\n• Environmental permits\n• Community acceptance"),
        ("Strategic Acquisitions", "• Distressed asset opportunities\n• Processing synergies\n• Vertical integration\n• Scale economics")
    ]
    
    x_positions = [Inches(1), Inches(4.7), Inches(8.4)]
    
    for i, (column_title, column_content) in enumerate(columns):
        # Column title
        col_title = slide.shapes.add_textbox(x_positions[i], Inches(2), Inches(3.5), Inches(0.6))
        col_title_frame = col_title.text_frame
        col_title_frame.text = column_title
        p = col_title_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        # Column content
        col_content = slide.shapes.add_textbox(x_positions[i], Inches(2.8), Inches(3.5), Inches(3))
        col_content_frame = col_content.text_frame
        col_content_frame.text = column_content
        col_content_frame.word_wrap = True
        
        for paragraph in col_content_frame.paragraphs:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = DARK_NAVY
            paragraph.alignment = PP_ALIGN.LEFT
    
    # Footer note
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.6))
    footer_frame = footer.text_frame
    footer_frame.text = "Pipeline: Multiple tenements under active review"
    p = footer_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

def add_three_engines_slide(prs):
    """Slide 7: Three Integrated Engines"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(0.7))
    title_frame = title.text_frame
    title_frame.text = "Platform Architecture"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(0.5))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Mining, R&D, and Manufacturing as integrated engines"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Three columns
    engines = [
        ("Mining & Exploration", "Fast Track to Revenue", "• Gold & silver production\n• Critical minerals exploration\n• Parallel development\n• Processing infrastructure"),
        ("R&D", "Critical Demand for Innovation", "• Leaching optimization\n• REE processing\n• Mining automation\n• Intellectual property"),
        ("Manufacturing", "Supply Chain Demand", "• Modular processing plants\n• Standardised deployment\n• Unit economics\n• Scalable operations")
    ]
    
    x_positions = [Inches(1), Inches(4.7), Inches(8.4)]
    
    for i, (engine_title, engine_subtitle, engine_content) in enumerate(engines):
        # Engine box
        engine_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_positions[i], Inches(2), Inches(3.5), Inches(4))
        engine_box.fill.solid()
        engine_box.fill.fore_color.rgb = LIGHT_GRAY
        engine_box.line.color.rgb = GOLD
        engine_box.line.width = Pt(3)
        
        # Engine title
        eng_title = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), Inches(2.2), Inches(3.3), Inches(0.6))
        eng_title_frame = eng_title.text_frame
        eng_title_frame.text = engine_title
        p = eng_title_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        # Engine subtitle
        eng_subtitle = slide.shapes.add_textbox(x_positions[i] + Inches(0.1), Inches(2.8), Inches(3.3), Inches(0.4))
        eng_subtitle_frame = eng_subtitle.text_frame
        eng_subtitle_frame.text = engine_subtitle
        p = eng_subtitle_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
        
        # Engine content
        eng_content = slide.shapes.add_textbox(x_positions[i] + Inches(0.2), Inches(3.4), Inches(3.1), Inches(2.4))
        eng_content_frame = eng_content.text_frame
        eng_content_frame.text = engine_content
        eng_content_frame.word_wrap = True
        
        for paragraph in eng_content_frame.paragraphs:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = DARK_NAVY
            paragraph.alignment = PP_ALIGN.LEFT
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    footer_frame = footer.text_frame
    footer_frame.text = "Integrated platform creates competitive moats"
    p = footer_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

def add_mining_exploration_slide(prs):
    """Slide 8: Mining & Exploration"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(0.7))
    title_frame = title.text_frame
    title_frame.text = "Mining & Exploration"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Callout
    callout = slide.shapes.add_textbox(Inches(8), Inches(1.2), Inches(4), Inches(0.8))
    callout_frame = callout.text_frame
    callout_frame.text = "12-18 months to revenue"
    p = callout_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Callout background
    callout_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(1), Inches(4.4), Inches(1.2))
    callout_bg.fill.solid()
    callout_bg.fill.fore_color.rgb = GOLD
    callout_bg.line.width = Pt(0)
    # Move callout text to front
    slide.shapes._spTree.remove(callout_bg._element)
    slide.shapes._spTree.insert(-1, callout_bg._element)
    
    # Five-step framework
    steps = [
        "Pilot Processing",
        "Toll Treatment", 
        "Dynamic Routing",
        "Fast to Cash",
        "Parallel Exploration"
    ]
    
    # Create arrow flow
    arrow_y = Inches(3.5)
    step_width = Inches(2)
    step_spacing = Inches(2.3)
    start_x = Inches(1.5)
    
    for i, step in enumerate(steps):
        x = start_x + (i * step_spacing)
        
        # Step box
        step_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, arrow_y, step_width, Inches(0.8))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = DARK_NAVY
        step_box.line.width = Pt(0)
        
        # Step text
        step_text = slide.shapes.add_textbox(x + Inches(0.1), arrow_y + Inches(0.15), step_width - Inches(0.2), Inches(0.5))
        step_text_frame = step_text.text_frame
        step_text_frame.text = step
        p = step_text_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow (except for last step)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_connector(1, x + step_width, arrow_y + Inches(0.4), x + step_spacing - Inches(0.05), arrow_y + Inches(0.4))
            arrow.line.color.rgb = GOLD
            arrow.line.width = Pt(3)
    
    # Additional content
    content = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    content_frame = content.text_frame
    content_frame.text = "Integrated approach reduces time-to-cash through processing access while building long-term exploration pipeline. Parallel development maximizes resource utilization and cash generation."
    content_frame.word_wrap = True
    p = content_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER

def add_simple_slide(prs, slide_num, title_text, content_items=None, subtitle_text=None, is_dark=False):
    """Helper function to create simpler slides"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    if is_dark:
        fill.fore_color.rgb = DARK_NAVY
        title_color = WHITE
        content_color = WHITE
    else:
        fill.fore_color.rgb = WHITE
        title_color = DARK_NAVY
        content_color = DARK_NAVY
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle if provided
    if subtitle_text:
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.333), Inches(0.5))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = subtitle_text
        p = subtitle_frame.paragraphs[0]
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
    
    # Content items if provided
    if content_items:
        start_y = Inches(2.2) if subtitle_text else Inches(1.5)
        content = slide.shapes.add_textbox(Inches(2), start_y, Inches(9.333), Inches(4))
        content_frame = content.text_frame
        content_frame.text = "\n\n".join(content_items)
        content_frame.word_wrap = True
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = content_color
            paragraph.alignment = PP_ALIGN.LEFT

def create_presentation():
    """Create the complete MineralX Strategic Fund Vision Deck"""
    prs = create_blank_presentation()
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Why Now
    add_why_now_slide(prs)
    
    # Slide 3: Industry Problem (Dark)
    add_industry_problem_slide(prs)
    
    # Slide 4: Platform
    add_platform_slide(prs)
    
    # Slide 5: Traction
    add_traction_slide(prs)
    
    # Slide 6: Growth Strategy
    add_growth_strategy_slide(prs)
    
    # Slide 7: Three Engines
    add_three_engines_slide(prs)
    
    # Slide 8: Mining & Exploration
    add_mining_exploration_slide(prs)
    
    # Slide 9: R&D
    add_simple_slide(prs, 9, "R&D", [
        "• Leaching Optimization: Advanced hydrometallurgical processes for complex ore bodies",
        "• REE Processing: Domestic rare earth element separation and purification",
        "• Mining Automation: Remote operations and process optimization systems",
        "• Process IP: Proprietary techniques for mineral extraction and processing",
        "",
        "R&D funded by operations + government grants. Targeting significant non-dilutive funding."
    ])
    
    # Slide 10: Manufacturing
    add_simple_slide(prs, 10, "Manufacturing", [
        "• Modularity: Pre-fabricated processing units for rapid deployment",
        "• Standardisation: Proven designs reduce engineering risk and cost", 
        "• Deployment Speed: 6-12 month installation vs 2-3 years traditional",
        "• Unit Economics: Optimized for Australian mineral characteristics and logistics"
    ])
    
    # Slide 11: Market Opportunity
    add_simple_slide(prs, 11, "Market Opportunity", [
        "• Gold: Global market $230B+, Australia 2nd largest producer",
        "• Critical Minerals: $770B by 2030, supply deficit growing",
        "• Processing Gap: Domestic capacity insufficient for demand",
        "• QLD Opportunity: Fragmented districts ripe for consolidation",
        "",
        "Structural supply-demand imbalance creates multi-decade growth opportunity"
    ])
    
    # Slide 12: Commodities & Consolidation
    add_simple_slide(prs, 12, "Commodities & Consolidation", [
        "• Gold & Silver: Immediate liquidity, proven metallurgy, 12-18 month path",
        "• Critical Minerals & REEs: Structural demand, longer timeline offset by gold cashflow", 
        "• Consolidation Strategy: Targeting fragmented districts, processing as lever",
        "• Acquisition Criteria: Proven mineralisation, fragmented ownership, no processing access, district-scale potential"
    ])
    
    # Slide 13: Capital Structure
    add_simple_slide(prs, 13, "Capital Structure", [
        "1. Operating Cash Flow → Self-funded growth from gold production",
        "2. NSR/Royalties → Asset-backed financing for expansion projects",
        "3. Project-Level Debt → Non-recourse debt for large infrastructure",
        "4. Strategic Equity → Partnership capital for platform scaling",
        "",
        "Diversified capital structure minimizes dilution while funding growth"
    ])
    
    # Slide 14: The Ask (Dark/Gold)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
    title_frame = title.text_frame
    title_frame.text = "The Opportunity"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Deal terms
    deal_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Inches(9.333), Inches(2.5))
    deal_box.fill.solid()
    deal_box.fill.fore_color.rgb = WHITE
    deal_box.line.color.rgb = GOLD
    deal_box.line.width = Pt(3)
    
    deal_text = slide.shapes.add_textbox(Inches(2.5), Inches(1.8), Inches(8.333), Inches(1.8))
    deal_text_frame = deal_text.text_frame
    deal_text_frame.text = "Raise: AUD $10,000,000 | Equity: 5% at $200M post-money valuation\nStructure: Two-tranche ($6M signing + $4M milestone)"
    p = deal_text_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Use of proceeds
    proceeds_text = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(2))
    proceeds_text_frame = proceeds_text.text_frame
    proceeds_text_frame.text = """Use of Proceeds:
• Processing plant & equipment: 40%    • Mining operations & working capital: 25%
• Exploration & resource development: 15%    • Acquisitions & consolidation: 15%    • Corporate & G&A: 5%

Investor Rights: Board observer rights, quarterly reporting, pro-rata rights on future rounds"""
    proceeds_text_frame.word_wrap = True
    
    for paragraph in proceeds_text_frame.paragraphs:
        paragraph.font.name = "Calibri"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Slide 15: Structural Tailwinds
    add_simple_slide(prs, 15, "Structural Tailwinds", [
        "• Exploration Grants: GEOVIC, Collaborative Drilling Program, Junior Minerals Exploration Incentive",
        "• Processing Incentives: Critical Minerals Facility loans, Regional Investment Corporation funding",
        "• R&D Funding: R&D Tax Incentive (43.5%), CRC-P grants, ARENA technology development",
        "• Strategic Alignment: Australia's Critical Minerals Strategy, supply chain security mandates"
    ], "Government co-funding opportunities")
    
    # Slide 16: Deployment Timeline
    add_simple_slide(prs, 16, "Deployment Timeline", [
        "Q1-Q4 2026: Initial Operations",
        "→ Pilot plant operational, toll treatment partnerships, early cash flow generation",
        "",
        "Q1-Q4 2027: District Consolidation", 
        "→ Strategic acquisitions, commercial scale processing, second plant deployment",
        "",
        "2028+: Multi-District Platform",
        "→ Multiple concurrent operations, standardised deployments, established market position"
    ])
    
    # Slide 17: Team
    add_simple_slide(prs, 17, "Leadership", [
        "Ayrton Mansi — Founder & CEO",
        "• Chemical engineering background with process design expertise",
        "• Hands-on mining & processing operational experience", 
        "• Proven property development and project management track record",
        "• Current operator of ML 30139 gold mining lease",
        "",
        "(Advisory board and key hires to be announced)"
    ])
    
    # Slide 18: Closing (Dark)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1))
    title_frame = title.text_frame
    title_frame.text = "MINERALX"
    p = title_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(0.6))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Integrated Mining & Processing Platform"
    p = subtitle_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Three bullet summary
    summary = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(2))
    summary_frame = summary.text_frame
    summary_frame.text = """• Processing infrastructure as the consolidation lever
• Near-term gold cashflow funds long-term platform value  
• First-mover in fragmented Australian mineral districts"""
    summary_frame.word_wrap = True
    
    for paragraph in summary_frame.paragraphs:
        paragraph.font.name = "Calibri"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.LEFT
    
    # Contact
    contact = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.333), Inches(0.6))
    contact_frame = contact.text_frame
    contact_frame.text = "Contact: Ayrton Mansi | ayrton@mineral-x.com.au"
    p = contact_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    footer_frame = footer.text_frame
    footer_frame.text = "Confidential | February 2026"
    p = footer_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    print("Creating MineralX Strategic Fund Vision Deck...")
    presentation = create_presentation()
    
    output_path = "/Users/ayrtonmansi/.openclaw/workspace/MineralX_Vision_Deck.pptx"
    presentation.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print("📊 18 slides created with professional design")
    print("🎯 Ready for strategic capital partners")